import io
import zipfile
from unittest.mock import AsyncMock, Mock, patch

import httpx
import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from kolega_code.agent.tool_backend.web_fetch.documents import DocumentConverter
from kolega_code.agent.tool_backend.web_fetch.extractors import extract_html, quality_score
from kolega_code.agent.tool_backend.web_fetch.pipeline import LocalWebContentPipeline, WebContentError
from kolega_code.agent.tool_backend.web_fetch.retrieval import (
    TEXT_MAX_BYTES,
    FetchedResource,
    RetrievalError,
    WebRetriever,
    normalize_url,
)


def _client_factory(handler):
    def factory(**kwargs):
        return httpx.AsyncClient(transport=httpx.MockTransport(handler), **kwargs)

    return factory


def test_normalize_url_repairs_model_variants() -> None:
    assert normalize_url("example.com/docs#part") == "https://example.com/docs"
    assert normalize_url("https:/example.com/a") == "https://example.com/a"
    assert normalize_url("<HTTP://EXAMPLE.COM>") == "http://example.com/"
    with pytest.raises(RetrievalError, match=r"http\(s\)"):
        normalize_url("ftp://example.com")
    with pytest.raises(RetrievalError, match="credentials"):
        normalize_url("https://user:pass@example.com")


@pytest.mark.asyncio
async def test_retriever_returns_metadata_and_decoded_body() -> None:
    def handler(request: httpx.Request) -> httpx.Response:
        assert request.headers["user-agent"]
        return httpx.Response(
            200,
            headers={"content-type": "text/plain; charset=iso-8859-1"},
            content="olá".encode("iso-8859-1"),
            request=request,
        )

    resource = await WebRetriever(_client_factory(handler)).fetch("example.com")
    assert resource.final_url == "https://example.com/"
    assert resource.content_type == "text/plain"
    assert resource.charset == "iso-8859-1"
    assert resource.body == "olá".encode("iso-8859-1")
    assert len(resource.attempts) == 1


@pytest.mark.asyncio
async def test_retriever_retries_transient_status_then_succeeds() -> None:
    calls = 0

    def handler(request: httpx.Request) -> httpx.Response:
        nonlocal calls
        calls += 1
        return httpx.Response(503 if calls == 1 else 200, text="temporary" if calls == 1 else "ready", request=request)

    with patch("kolega_code.agent.tool_backend.web_fetch.retrieval.asyncio.sleep", new=AsyncMock()):
        resource = await WebRetriever(_client_factory(handler)).fetch("https://example.com")
    assert resource.body == b"ready"
    assert calls == 2
    assert [attempt.status_code for attempt in resource.attempts] == [503, 200]


@pytest.mark.asyncio
async def test_retriever_rejects_declared_oversized_text() -> None:
    def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(
            200,
            headers={"content-length": str(TEXT_MAX_BYTES + 1), "content-type": "text/plain"},
            request=request,
        )

    with pytest.raises(RetrievalError, match="too large"):
        await WebRetriever(_client_factory(handler)).fetch("https://example.com/large")


@pytest.mark.asyncio
async def test_pipeline_handles_json_without_html_extractor() -> None:
    retriever = Mock()
    retriever.fetch = AsyncMock(
        return_value=FetchedResource(
            "https://api.example/",
            "https://api.example/",
            200,
            "application/json",
            "utf-8",
            None,
            b'{"answer":42}',
        )
    )
    pipeline = LocalWebContentPipeline(retriever=retriever)
    result = await pipeline.load("https://api.example")
    assert result.method == "json"
    assert result.content == '{\n  "answer": 42\n}'


@pytest.mark.asyncio
async def test_html_preference_falls_back_when_selected_extractor_is_empty() -> None:
    html = "<html><body><main><h1>Title</h1><p>" + ("Useful article text. " * 30) + "</p></main></body></html>"
    with patch(
        "kolega_code.agent.tool_backend.web_fetch.extractors.ReadabilityExtractor.extract",
        return_value="",
    ):
        result = await extract_html(html, "https://example.com", preference="readability")
    assert result.content
    assert result.method != "readability"
    assert result.attempts[0].name == "readability"


@pytest.mark.asyncio
async def test_spa_shell_is_detected_without_browser_fallback() -> None:
    html = """<html><body><div id="root"></div><noscript>Please enable JavaScript</noscript>
    <script src="runtime.js"></script><script src="react.js"></script><script src="app.js"></script></body></html>"""
    result = await extract_html(html, "https://spa.example")
    assert result.spa_detected is True
    assert any("does not run a browser" in warning for warning in result.warnings)


def test_quality_gate_accepts_legitimate_short_page() -> None:
    html = "<html><body><main><p>Small but complete answer.</p></main></body></html>"
    score, usable = quality_score("Small but complete answer with useful context for the reader.", html)
    assert usable is True
    assert score > 8


@pytest.mark.asyncio
async def test_markitdown_converts_xlsx_locally() -> None:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Facts"
    sheet.append(["Name", "Value"])
    sheet.append(["answer", 42])
    buffer = io.BytesIO()
    workbook.save(buffer)

    result = await DocumentConverter().convert(buffer.getvalue(), ".xlsx", "https://example.com/facts.xlsx")
    assert result.method == "markitdown:xlsx"
    assert "answer" in result.content and "42" in result.content


@pytest.mark.asyncio
async def test_markitdown_converts_docx_locally() -> None:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
              <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
              <Default Extension="xml" ContentType="application/xml"/>
              <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
              <w:body><w:p><w:r><w:t>Reliable DOCX fact</w:t></w:r></w:p><w:sectPr/></w:body>
            </w:document>""",
        )

    result = await DocumentConverter().convert(buffer.getvalue(), ".docx", "https://example.com/facts.docx")
    assert "Reliable DOCX fact" in result.content


@pytest.mark.asyncio
async def test_markitdown_converts_pptx_locally() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    text_box.text_frame.text = "Reliable PPTX fact"
    buffer = io.BytesIO()
    presentation.save(buffer)

    result = await DocumentConverter().convert(buffer.getvalue(), ".pptx", "https://example.com/facts.pptx")
    assert "Reliable PPTX fact" in result.content


@pytest.mark.asyncio
async def test_markitdown_reports_image_only_pdf() -> None:
    converter = DocumentConverter()
    fake = Mock()
    fake.convert_stream.return_value = Mock(text_content="")
    converter._converter = fake
    with pytest.raises(Exception, match="scanned or image-only"):
        await converter.convert(b"%PDF-empty", ".pdf", "https://example.com/scan.pdf")


@pytest.mark.asyncio
async def test_pipeline_reports_unsupported_binary() -> None:
    retriever = Mock()
    retriever.fetch = AsyncMock(
        return_value=FetchedResource(
            "https://example.com/file.zip",
            "https://example.com/file.zip",
            200,
            "application/zip",
            None,
            None,
            b"not-a-supported-archive",
        )
    )
    pipeline = LocalWebContentPipeline(retriever=retriever)
    with pytest.raises(WebContentError, match="Unsupported content type"):
        await pipeline.load("https://example.com/file.zip")


@pytest.mark.asyncio
async def test_pipeline_reports_unrecognized_octet_stream_as_binary() -> None:
    retriever = Mock()
    retriever.fetch = AsyncMock(
        return_value=FetchedResource(
            "https://example.com/blob",
            "https://example.com/blob",
            200,
            "application/octet-stream",
            None,
            None,
            b"\x00\x01\x02opaque-binary",
        )
    )
    with pytest.raises(WebContentError, match="binary"):
        await LocalWebContentPipeline(retriever=retriever).load("https://example.com/blob")


@pytest.mark.asyncio
async def test_pipeline_reports_legacy_office_format_precisely() -> None:
    retriever = Mock()
    retriever.fetch = AsyncMock(
        return_value=FetchedResource(
            "https://example.com/old.doc",
            "https://example.com/old.doc",
            200,
            "application/msword",
            None,
            None,
            b"legacy-doc-bytes",
        )
    )
    with pytest.raises(WebContentError, match="Legacy DOC"):
        await LocalWebContentPipeline(retriever=retriever).load("https://example.com/old.doc")
